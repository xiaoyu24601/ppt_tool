"""
PPT 工具箱 — AI 生成 · 美化 · 去水印 · 修复
管理密码 + 激活码 + 免 API 直接使用
"""
import streamlit as st
import hashlib, hmac, secrets, json, io, os, zipfile, tempfile, shutil, uuid, re, base64
from pathlib import Path
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import xml.etree.ElementTree as ET
import requests

st.set_page_config(page_title="PPT 工具箱", page_icon="🪄", layout="wide")

# ==================== 配置 ====================

# DeepSeek API Key（你的密钥，内置在代码中）
DEEPSEEK_API_KEY = "sk-7e6523558c3a457c93b60b2f87a9155c"

# 管理密码 — 改成你自己的
ADMIN_PASSWORD = "24601"

# 激活码密钥盐
ACTIVATION_SALT = "ppt-tool-salt-2026"

# 预生成的有效激活码（用 generate_code() 生成）
# 格式: {"code": "激活码", "used": False, "created": "日期", "note": "备注"}
VALID_CODES = [
    {"code": "PPT-VIP-0001", "used": False, "created": "2026-05-13", "note": "测试码"},
    {"code": "PPT-VIP-0002", "used": False, "created": "2026-05-13", "note": "测试码"},
    {"code": "PPT-VIP-0003", "used": False, "created": "2026-05-13", "note": "测试码"},
]

# ==================== 激活码逻辑 ====================

def generate_code(prefix="PPT", index=1):
    """生成新激活码"""
    raw = f"{ACTIVATION_SALT}-{prefix}-{index}-{secrets.token_hex(4)}"
    hash_val = hashlib.sha256(raw.encode()).hexdigest()[:8].upper()
    return f"{prefix}-{hash_val}"

def validate_code(code):
    """验证激活码是否有效且未使用"""
    for entry in VALID_CODES:
        if entry["code"] == code and not entry["used"]:
            return entry
    return None

def mark_code_used(code):
    """标记激活码已使用"""
    for entry in VALID_CODES:
        if entry["code"] == code:
            entry["used"] = True
            entry["used_at"] = datetime.now().isoformat()
            return True
    return False

# ==================== 会话状态 ====================

if "ai_enabled" not in st.session_state:
    st.session_state.ai_enabled = False
if "auth_type" not in st.session_state:
    st.session_state.auth_type = None  # "admin" or "activation"
if "code_used" not in st.session_state:
    st.session_state.code_used = None

# ==================== 主题配置 ====================

THEMES = {
    "专业商务": {"primary": "#2C3E50", "secondary": "#3498DB", "accent": "#2ECC71", "font": "Arial"},
    "学术答辩": {"primary": "#003366", "secondary": "#C0392B", "accent": "#F39C12", "font": "Times New Roman"},
    "创业路演": {"primary": "#E74C3C", "secondary": "#F39C12", "accent": "#9B59B6", "font": "Arial"},
    "科技互联网": {"primary": "#6C63FF", "secondary": "#00D4FF", "accent": "#FF6B6B", "font": "Helvetica"},
    "教育培训": {"primary": "#27AE60", "secondary": "#2ECC71", "accent": "#3498DB", "font": "Arial"},
    "简约极简": {"primary": "#333333", "secondary": "#666666", "accent": "#999999", "font": "Helvetica"},
    "清新自然": {"primary": "#16A085", "secondary": "#27AE60", "accent": "#F1C40F", "font": "Arial"},
    "暗夜模式": {"primary": "#FFFFFF", "secondary": "#3498DB", "accent": "#E74C3C", "font": "Arial", "dark": True},
    "创意设计": {"primary": "#9B59B6", "secondary": "#E91E63", "accent": "#FF9800", "font": "Comic Sans MS"},
    "政府汇报": {"primary": "#8B0000", "secondary": "#B8860B", "accent": "#006400", "font": "SimHei"},
}

OUTLINE_LIBRARY = {
    "专业商务": ["项目背景", "市场分析", "解决方案", "竞争优势", "执行计划", "预期成果"],
    "学术答辩": ["研究背景", "文献综述", "研究方法", "实验结果", "分析讨论", "结论展望"],
    "创业路演": ["痛点分析", "解决方案", "市场规模", "商业模式", "团队介绍", "融资计划"],
    "科技互联网": ["行业趋势", "技术架构", "产品功能", "用户场景", "数据指标", "未来路线"],
    "教育培训": ["课程目标", "核心概念", "案例讲解", "实操演练", "常见误区", "总结回顾"],
    "简约极简": ["背景", "问题", "方案", "优势", "计划", "总结"],
    "清新自然": ["引入", "现状", "对策", "案例", "成果", "展望"],
    "暗夜模式": ["背景", "挑战", "策略", "执行", "数据", "下一步"],
    "创意设计": ["灵感来源", "设计理念", "视觉方案", "细节展示", "应用场景", "迭代方向"],
    "政府汇报": ["工作概述", "重点任务", "进展成效", "问题分析", "改进措施", "下步计划"],
}

AI_RESEARCH_PROMPT = """你是一个专业的 PPT 策划顾问。用户给了你一个主题，你需要像麦肯锡顾问一样，先研究这个主题的市场上顶级 PPT 应该包含什么内容。

主题：{topic}
场景：{theme}

请分析并返回 JSON：
{{
  "market_context": "这个主题的市场背景和当前热点（1-2句话）",
  "target_audience": "目标听众是谁",
  "key_angles": ["角度1", "角度2", "角度3"],
  "data_sources": ["数据来源思路1", "数据来源思路2"],
  "outline": [
    {{"section": "章节标题", "focus": "这一章要讲什么核心内容", "keywords": ["关键词1", "关键词2"]}}
  ]
}}

注意：返回的 outline 要包含封面页作为第一项，一共{slide_count}页左右。每页都要有明确的商业/学术价值。"""

AI_KIMI_PROMPT = """你是一个顶级 PPT 策划师。用户要把下面的内容交给 Kimi 来生成 PPT。你需要写一份**直接复制粘贴给 Kimi 的 Markdown 提示词**。

## 研究报告
{research_result}

## 用户原始主题
{user_prompt}

## 场景：{theme}  |  目标页数：约{slide_count}页

## 核心任务
写一份可以直接粘贴给 Kimi 的 PPT 生成指令。必须包含：

### 1. 角色设定
告诉Kimi它扮演什么角色（如"资深商业分析师""麦肯锡顾问"等），让它进入专业状态。

### 2. 整体视觉风格
- 配色方案建议（主色/辅色/强调色）
- 字体风格建议
- 整体氛围（专业/创新/温暖/科技感等）

### 3. 逐页详细内容（最重要！）
每页都要写明：
- 页面标题
- **布局建议**（关键！不要全用要点列表，要多样化）：
  - 有的用左右分栏对比
  - 有的用大数字突出关键指标
  - 有的用表格对比
  - 有的用时间线
  - 有的用图标+文字卡片
  - 有的用引用/金句强调
- 具体文字内容（不是提纲，是可用的完整文案）
- 可以建议配图方向（如"这里放一张市场规模趋势图"）

### 4. 格式要求
- 直接用 Markdown 格式
- 中文输出
- 内容具体，有数据有案例，拒绝空洞套话
- 每页布局必须有变化，这是专业PPT的基本要求

直接输出 Markdown 文本（不要用代码块包裹，不要加 json 之类的标记）："""

AI_PPTX_PROMPT = """你是一个 PPT 数据格式化助手。将以下 Kimi 提示词内容转换为 PPTX 生成用的 JSON。

## Kimi 提示词内容
{kimi_markdown}

## 页数限制
约{slide_count}页

输出严格的 JSON（不要 markdown）：
{{
  "title": "PPT 主标题",
  "subtitle": "副标题",
  "slides": [
    {{
      "title": "页面标题",
      "layout": "bullets",
      "bullets": ["要点1（要有实质内容）", "要点2", "要点3"],
      "notes": "演讲备注"
    }},
    {{
      "title": "关键数据",
      "layout": "big_number",
      "big_number": {{"value": "500亿", "label": "2025市场规模"}},
      "bullets": ["补充说明"],
      "notes": ""
    }},
    {{
      "title": "方案对比",
      "layout": "comparison",
      "columns": [
        {{"heading": "方案A", "items": ["特点1", "特点2"]}},
        {{"heading": "方案B", "items": ["特点1", "特点2"]}}
      ],
      "notes": ""
    }}
  ]
}}

可用 layout 类型及所需字段：
- title_slide: 不需要额外字段，作为章节封面
- bullets: 需要 "bullets" 字段（字符串数组，3-5条）
- two_column: 需要 "columns" 字段（2列，每列有 heading + items）
- big_number: 需要 "big_number" 字段 {{"value": "", "label": ""}}，可选 "bullets"
- comparison: 需要 "columns" 字段（2-3列，每列有 heading + items）
- quote: 需要 "quote_text" 和 "quote_author" 字段
- timeline: 需要 "timeline_items" 字段 [{{"date": "", "event": ""}}]
- summary: 需要 "summary_items" 字段（字符串数组，3-4条结论）

重要：每页必须指定 layout，且要多样化，不要全部用 bullets！"""

# ==================== 注册 XML 命名空间 ====================

for p, u in {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}.items():
    ET.register_namespace(p, u)

OUTPUT_DIR = Path(__file__).parent / "outputs"
OUTPUT_DIR.mkdir(exist_ok=True)


# ==================== AI PPT 生成 ====================

def call_deepseek(messages, max_tokens=4096):
    """调用 DeepSeek API"""
    resp = requests.post(
        "https://api.deepseek.com/chat/completions",
        headers={"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"},
        json={"model": "deepseek-chat", "messages": messages, "temperature": 0.7, "max_tokens": max_tokens},
        timeout=120,
    )
    if resp.status_code != 200:
        raise Exception(f"API 错误 ({resp.status_code}): {resp.text[:200]}")
    return resp.json()["choices"][0]["message"]["content"]


def parse_json(raw):
    raw = raw.strip()
    if raw.startswith("```"):
        raw = re.sub(r'^```\w*\n', '', raw)
        raw = re.sub(r'\n```$', '', raw)
    return json.loads(raw)


def generate_ppt_content(prompt, theme, slide_count):
    """三步生成法：研究 → Kimi Markdown（主输出）→ PPTX JSON（草稿）"""
    # 第一步：研究大纲
    research_msg = [
        {"role": "user", "content": AI_RESEARCH_PROMPT.format(topic=prompt, theme=theme, slide_count=slide_count)}
    ]
    raw_research = call_deepseek(research_msg, 2048)
    research = parse_json(raw_research)

    # 第二步：生成 Kimi 优化的 Markdown 提示词（主输出）
    kimi_msg = [
        {"role": "user", "content": AI_KIMI_PROMPT.format(
            user_prompt=prompt,
            research_result=json.dumps(research, ensure_ascii=False, indent=2),
            theme=theme,
            slide_count=slide_count,
        )}
    ]
    kimi_markdown = call_deepseek(kimi_msg, 4096)

    # 第三步：基于 Markdown 生成 PPTX JSON（作为草稿下载）
    pptx_msg = [
        {"role": "user", "content": AI_PPTX_PROMPT.format(
            kimi_markdown=kimi_markdown[:3000],
            slide_count=slide_count,
        )}
    ]
    raw_content = call_deepseek(pptx_msg, 2048)
    content = parse_json(raw_content)
    return content, research, kimi_markdown


def build_pptx(data, theme_name, slide_count):
    """多布局 PPTX 构建器 — 支持 8 种布局类型"""
    t = THEMES.get(theme_name, THEMES["专业商务"])
    is_dark = t.get("dark", False)
    pc, sc, ac = t["primary"], t["secondary"], t["accent"]
    bg = "#1E1E1E" if is_dark else "#FFFFFF"
    tc = "#FFFFFF" if is_dark else "#333333"
    font = t["font"]

    def h(c):
        c = c.lstrip("#")
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_text = data.get("title", "未命名")
    subtitle_text = data.get("subtitle", "")
    slides = data.get("slides", [])
    if not slides:
        slides = [{"title": title_text, "layout": "bullets", "bullets": ["请查看详细内容"], "notes": ""}]

    def add_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = h(bg)

    def add_pn(slide, num):
        tb = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        p.text = str(num); p.font.size = Pt(11)
        p.font.color.rgb = h("#AAAAAA"); p.alignment = PP_ALIGN.RIGHT

    def add_title(slide, title):
        bar = slide.shapes.add_shape(1, Inches(0.7), Inches(0.45), Inches(0.07), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = h(sc); bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(0.35), Inches(10.5), Inches(0.7))
        tp = tb.text_frame.paragraphs[0]
        tp.text = title; tp.font.size = Pt(30); tp.font.bold = True
        tp.font.color.rgb = h(pc); tp.font.name = font
        ln = slide.shapes.add_shape(1, Inches(1.0), Inches(1.15), Inches(2.0), Inches(0.022))
        ln.fill.solid(); ln.fill.fore_color.rgb = h(sc); ln.line.fill.background()

    def add_notes(slide, data_slide):
        if data_slide.get("notes"):
            try:
                slide.notes_slide.notes_text_frame.text = data_slide["notes"]
            except:
                pass

    # === 封面 ===
    sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl)
    deco = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.13), Inches(7.5))
    deco.fill.solid(); deco.fill.fore_color.rgb = h(sc); deco.line.fill.background()
    deco2 = sl.shapes.add_shape(1, Inches(8.5), Inches(0), Inches(4.833), Inches(0.07))
    deco2.fill.solid(); deco2.fill.fore_color.rgb = h(ac); deco2.line.fill.background()
    tb = sl.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.5), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title_text[:60]; p.font.size = Pt(52); p.font.bold = True
    p.font.color.rgb = h(pc); p.font.name = font
    if subtitle_text:
        tb2 = sl.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(10.5), Inches(1.2))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle_text[:120]; p2.font.size = Pt(20)
        p2.font.color.rgb = h("#888888"); p2.font.name = font
    tb3 = sl.shapes.add_textbox(Inches(1.2), Inches(6.4), Inches(10.5), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = f"{datetime.now().strftime('%Y.%m.%d')}  |  {theme_name}"
    p3.font.size = Pt(12); p3.font.color.rgb = h("#AAAAAA")

    # === 内容页 ===
    for i, sd in enumerate(slides[:slide_count]):
        sl = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(sl)
        layout = sd.get("layout", "bullets")
        stitle = sd.get("title", f"第{i+1}页")
        card_c = ["#F5F5F5" if not is_dark else "#2A2A2A",
                  "#FAFAFA" if not is_dark else "#252525"]

        if layout == "title_slide":
            deco = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.13), Inches(7.5))
            deco.fill.solid(); deco.fill.fore_color.rgb = h(sc); deco.line.fill.background()
            tb = sl.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
            p = tb.text_frame.paragraphs[0]
            p.text = stitle; p.font.size = Pt(48); p.font.bold = True
            p.font.color.rgb = h(pc); p.font.name = font
            bullets = sd.get("bullets", [])
            if bullets:
                tb2 = sl.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(0.8))
                p2 = tb2.text_frame.paragraphs[0]
                p2.text = bullets[0]; p2.font.size = Pt(18)
                p2.font.color.rgb = h("#999999"); p2.font.name = font
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "bullets":
            add_title(sl, stitle)
            bullets = sd.get("bullets", ["暂无内容"])
            colors3 = [pc, sc, ac]
            for j, bullet in enumerate(bullets):
                y = Inches(1.55 + j * 1.05)
                circle = sl.shapes.add_shape(9, Inches(1.1), y + Inches(0.03), Inches(0.42), Inches(0.42))
                circle.fill.solid(); circle.fill.fore_color.rgb = h(colors3[j % 3]); circle.line.fill.background()
                cp = circle.text_frame.paragraphs[0]
                cp.text = str(j + 1); cp.font.size = Pt(13); cp.font.bold = True
                cp.font.color.rgb = h("#FFFFFF"); cp.alignment = PP_ALIGN.CENTER
                card = sl.shapes.add_shape(5, Inches(1.75), y, Inches(10.2), Inches(0.78))
                card.fill.solid()
                try: card.fill.fore_color.rgb = h(card_c[j % 2])
                except: pass
                card.line.fill.background()
                tb = sl.shapes.add_textbox(Inches(2.05), y + Inches(0.05), Inches(9.7), Inches(0.68))
                tp = tb.text_frame.paragraphs[0]
                tp.text = bullet[:110]; tp.font.size = Pt(16)
                tp.font.color.rgb = h(tc); tp.font.name = font
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "two_column":
            add_title(sl, stitle)
            cols = sd.get("columns", [{"heading": "左栏", "items": ["—"]}, {"heading": "右栏", "items": ["—"]}])
            for ci, col in enumerate(cols[:2]):
                x = Inches(1.0 + ci * 6.1)
                tb = sl.shapes.add_textbox(x, Inches(1.5), Inches(5.5), Inches(0.5))
                tp = tb.text_frame.paragraphs[0]
                tp.text = col.get("heading", ""); tp.font.size = Pt(20); tp.font.bold = True
                tp.font.color.rgb = h(sc); tp.font.name = font
                ln = sl.shapes.add_shape(1, x, Inches(2.05), Inches(1.3), Inches(0.02))
                ln.fill.solid(); ln.fill.fore_color.rgb = h(sc); ln.line.fill.background()
                for jj, item in enumerate(col.get("items", [])):
                    iy = Inches(2.3 + jj * 0.55)
                    tb2 = sl.shapes.add_textbox(x + Inches(0.15), iy, Inches(5.2), Inches(0.48))
                    tp2 = tb2.text_frame.paragraphs[0]
                    tp2.text = f"▸ {item[:80]}"; tp2.font.size = Pt(15)
                    tp2.font.color.rgb = h(tc); tp2.font.name = font
            if len(cols) >= 2:
                div = sl.shapes.add_shape(1, Inches(6.6), Inches(1.5), Inches(0.013), Inches(5.2))
                div.fill.solid(); div.fill.fore_color.rgb = h("#DDDDDD"); div.line.fill.background()
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "big_number":
            add_title(sl, stitle)
            bn = sd.get("big_number", {"value": "—", "label": ""})
            tb = sl.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5))
            tp = tb.text_frame.paragraphs[0]
            tp.text = str(bn.get("value", "—")); tp.font.size = Pt(80); tp.font.bold = True
            tp.font.color.rgb = h(sc); tp.font.name = font; tp.alignment = PP_ALIGN.CENTER
            label = bn.get("label", "")
            if label:
                tb2 = sl.shapes.add_textbox(Inches(2.0), Inches(4.6), Inches(9.3), Inches(0.7))
                tp2 = tb2.text_frame.paragraphs[0]
                tp2.text = label; tp2.font.size = Pt(22)
                tp2.font.color.rgb = h(tc); tp2.font.name = font; tp2.alignment = PP_ALIGN.CENTER
            for j, bullet in enumerate(sd.get("bullets", [])[:2]):
                tb3 = sl.shapes.add_textbox(Inches(2.5), Inches(5.4 + j * 0.42), Inches(8.3), Inches(0.38))
                tp3 = tb3.text_frame.paragraphs[0]
                tp3.text = f"• {bullet[:100]}"; tp3.font.size = Pt(14)
                tp3.font.color.rgb = h("#888888"); tp3.font.name = font; tp3.alignment = PP_ALIGN.CENTER
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "comparison":
            add_title(sl, stitle)
            cols = sd.get("columns", [])
            if not cols:
                cols = [{"heading": "A", "items": ["—"]}, {"heading": "B", "items": ["—"]}]
            n = min(len(cols), 3)
            cw = 10.5 / n
            for ci, col in enumerate(cols[:3]):
                x = Inches(1.3 + ci * (cw + 0.25))
                hdr = sl.shapes.add_shape(1, x, Inches(1.5), Inches(cw), Inches(0.52))
                hdr.fill.solid()
                hdr.fill.fore_color.rgb = h([pc, sc, ac][ci % 3])
                hdr.line.fill.background()
                hp = hdr.text_frame.paragraphs[0]
                hp.text = col.get("heading", ""); hp.font.size = Pt(17); hp.font.bold = True
                hp.font.color.rgb = h("#FFFFFF"); hp.alignment = PP_ALIGN.CENTER
                for jj, item in enumerate(col.get("items", [])):
                    iy = Inches(2.25 + jj * 0.52)
                    card = sl.shapes.add_shape(5, x, iy, Inches(cw), Inches(0.43))
                    card.fill.solid()
                    try: card.fill.fore_color.rgb = h(card_c[jj % 2])
                    except: pass
                    card.line.fill.background()
                    ctb = sl.shapes.add_textbox(x + Inches(0.12), iy + Inches(0.04), Inches(cw - 0.24), Inches(0.35))
                    ctp = ctb.text_frame.paragraphs[0]
                    ctp.text = item[:60]; ctp.font.size = Pt(13)
                    ctp.font.color.rgb = h(tc); ctp.font.name = font; ctp.alignment = PP_ALIGN.CENTER
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "quote":
            add_title(sl, stitle)
            qm = sl.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(2), Inches(2))
            qp = qm.text_frame.paragraphs[0]
            qp.text = "❭"; qp.font.size = Pt(72)
            qp.font.color.rgb = h(sc); qp.font.name = font
            qt = sd.get("quote_text", stitle)
            tb = sl.shapes.add_textbox(Inches(2.5), Inches(2.2), Inches(9.5), Inches(3))
            tp = tb.text_frame.paragraphs[0]
            tp.text = qt[:200]; tp.font.size = Pt(28); tp.font.italic = True
            tp.font.color.rgb = h(pc); tp.font.name = font
            qa = sd.get("quote_author", "")
            if qa:
                tb2 = sl.shapes.add_textbox(Inches(2.5), Inches(5.5), Inches(9.5), Inches(0.5))
                tp2 = tb2.text_frame.paragraphs[0]
                tp2.text = f"—— {qa}"; tp2.font.size = Pt(16)
                tp2.font.color.rgb = h("#999999"); tp2.font.name = font
                tp2.alignment = PP_ALIGN.RIGHT
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "timeline":
            add_title(sl, stitle)
            items = sd.get("timeline_items", [])
            if not items:
                items = [{"date": f"阶段{j+1}", "event": b}
                         for j, b in enumerate(sd.get("bullets", ["暂无内容"]))]
            ly = Inches(3.5)
            hline = sl.shapes.add_shape(1, Inches(1.2), ly, Inches(10.9), Inches(0.028))
            hline.fill.solid(); hline.fill.fore_color.rgb = h(sc); hline.line.fill.background()
            n = len(items)
            for j, item in enumerate(items[:6]):
                x = Inches(1.2 + j * (10.5 / max(n - 1, 1)))
                dot = sl.shapes.add_shape(9, x - Inches(0.09), ly - Inches(0.09), Inches(0.2), Inches(0.2))
                dot.fill.solid(); dot.fill.fore_color.rgb = h(sc); dot.line.fill.background()
                dtb = sl.shapes.add_textbox(x - Inches(0.55), ly - Inches(0.8), Inches(1.3), Inches(0.38))
                dp = dtb.text_frame.paragraphs[0]
                dp.text = item.get("date", ""); dp.font.size = Pt(13); dp.font.bold = True
                dp.font.color.rgb = h(sc); dp.font.name = font; dp.alignment = PP_ALIGN.CENTER
                etb = sl.shapes.add_textbox(x - Inches(0.65), ly + Inches(0.28), Inches(1.5), Inches(2.0))
                etb.text_frame.word_wrap = True
                ep = etb.text_frame.paragraphs[0]
                ep.text = item.get("event", "")[:60]; ep.font.size = Pt(11)
                ep.font.color.rgb = h(tc); ep.font.name = font; ep.alignment = PP_ALIGN.CENTER
            add_pn(sl, i + 1); add_notes(sl, sd)

        elif layout == "summary":
            add_title(sl, stitle)
            items = sd.get("summary_items", sd.get("bullets", ["谢谢"]))
            box_c = [pc, sc, ac]
            for j, item in enumerate(items[:4]):
                x = Inches(1.0 + (j % 2) * 6.1)
                y = Inches(1.55 + (j // 2) * 2.6)
                num = sl.shapes.add_textbox(x, y, Inches(0.75), Inches(0.75))
                np = num.text_frame.paragraphs[0]
                np.text = f"0{j+1}"; np.font.size = Pt(36); np.font.bold = True
                np.font.color.rgb = h(box_c[j % 3])
                card = sl.shapes.add_shape(5, x + Inches(0.65), y + Inches(0.08), Inches(4.85), Inches(0.85))
                card.fill.solid()
                try: card.fill.fore_color.rgb = h(card_c[0])
                except: pass
                card.line.fill.background()
                ctb = sl.shapes.add_textbox(x + Inches(0.85), y + Inches(0.12), Inches(4.45), Inches(0.75))
                ctp = ctb.text_frame.paragraphs[0]
                ctp.text = item[:110]; ctp.font.size = Pt(15)
                ctp.font.color.rgb = h(tc); ctp.font.name = font
            add_pn(sl, i + 1); add_notes(sl, sd)

        else:
            # fallback
            add_title(sl, stitle)
            for j, bullet in enumerate(sd.get("bullets", ["暂无内容"])):
                y = Inches(1.55 + j * 0.65)
                tb = sl.shapes.add_textbox(Inches(1.3), y, Inches(11), Inches(0.5))
                tp = tb.text_frame.paragraphs[0]
                tp.text = f"▸ {bullet[:120]}"; tp.font.size = Pt(18)
                tp.font.color.rgb = h(tc); tp.font.name = font
            add_pn(sl, i + 1); add_notes(sl, sd)

    output = io.BytesIO()
    prs.save(output); output.seek(0)
    return output


# ==================== 通用去水印 ====================

COMMON_WATERMARKS = [
    "made with gamma", "gamma.app", "gamma", "canva", "wps", "slidesgo",
    "slideshare", "prezi", "beautiful.ai", "slidebean", "powtoon",
    "visme", "piktochart", "genially", "googleslides", "keynote",
    "watermark", "sample", "preview", "draft", "confidential", "demo", "trial",
]


def remove_watermark(file_bytes, filename, target=None):
    temp_dir, extract_dir = tempfile.mkdtemp(), None
    try:
        ip = os.path.join(temp_dir, filename)
        op = os.path.join(temp_dir, f"clean_{filename}")
        with open(ip, 'wb') as f: f.write(file_bytes)
        extract_dir = temp_dir + "_ex"
        with zipfile.ZipFile(ip, 'r') as z: z.extractall(extract_dir)
        keywords = [target.lower()] if target else COMMON_WATERMARKS
        found, removed = set(), 0
        for root_dir, dirs, files in os.walk(extract_dir):
            for file in files:
                if not file.endswith('.xml'): continue
                fp = os.path.join(root_dir, file)
                try:
                    tree = ET.parse(fp); root = tree.getroot(); changed = False
                    for elem in root.iter():
                        if elem.text:
                            for kw in keywords:
                                if kw in elem.text.lower():
                                    found.add(kw); elem.text = ''; changed = True; removed += 1; break
                        for an, av in list(elem.attrib.items()):
                            if av:
                                for kw in keywords:
                                    if kw in str(av).lower():
                                        found.add(kw); del elem.attrib[an]; changed = True; removed += 1; break
                    if changed:
                        tree.write(fp, xml_declaration=True, encoding='UTF-8')
                except ET.ParseError: continue
        with zipfile.ZipFile(op, 'w', zipfile.ZIP_DEFLATED) as zout:
            for root_dir, dirs, files in os.walk(extract_dir):
                for file in files:
                    fp = os.path.join(root_dir, file)
                    zout.write(fp, os.path.relpath(fp, extract_dir))
        with open(op, 'rb') as f: return f.read(), removed, list(found)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
        if extract_dir: shutil.rmtree(extract_dir, ignore_errors=True)


# ==================== 布局修复 ====================

def fix_layout(file_bytes, filename):
    temp_dir, extract_dir = tempfile.mkdtemp(), None
    try:
        ip = os.path.join(temp_dir, filename)
        op = os.path.join(temp_dir, f"fixed_{filename}")
        with open(ip, 'wb') as f: f.write(file_bytes)
        extract_dir = temp_dir + "_ex"
        with zipfile.ZipFile(ip, 'r') as z: z.extractall(extract_dir)
        fixed = 0
        for root_dir, dirs, files in os.walk(extract_dir):
            for file in files:
                if not file.endswith('.xml'): continue
                fp = os.path.join(root_dir, file)
                try:
                    tree = ET.parse(fp); root = tree.getroot(); changed = False
                    for elem in root.iter():
                        tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                        if tag in ('rPr', 'defRPr'):
                            for child in elem:
                                ct = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                                if ct == 'latin' and not child.get('typeface', ''):
                                    child.set('typeface', 'Arial'); changed = True
                                elif ct == 'ea' and not child.get('typeface', ''):
                                    child.set('typeface', '微软雅黑'); changed = True
                    if changed: tree.write(fp, xml_declaration=True, encoding='UTF-8'); fixed += 1
                except ET.ParseError: continue
        with zipfile.ZipFile(op, 'w', zipfile.ZIP_DEFLATED) as zout:
            for root_dir, dirs, files in os.walk(extract_dir):
                for file in files:
                    fp = os.path.join(root_dir, file)
                    zout.write(fp, os.path.relpath(fp, extract_dir))
        with open(op, 'rb') as f: return f.read(), fixed
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
        if extract_dir: shutil.rmtree(extract_dir, ignore_errors=True)


# ==================== PPT 美化 ====================

def beautify_pptx(file_bytes, filename, theme_name, intensity):
    t = THEMES.get(theme_name, THEMES["专业商务"])
    is_dark = t.get("dark", False)
    pc, sc = t["primary"], t["secondary"]
    bg = "#1E1E1E" if is_dark else "#FFFFFF"
    tc = "#FFFFFF" if is_dark else "#333333"
    font = t["font"]

    def h(c):
        c = c.lstrip("#")
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

    prs = Presentation(io.BytesIO(file_bytes))
    changes = 0

    for master in prs.slide_masters:
        try:
            master.background.fill.solid(); master.background.fill.fore_color.rgb = h(bg)
        except: pass

    for slide in prs.slides:
        try:
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = h(bg)
        except: pass
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font
                        if intensity != "轻度":
                            if run.font.size and run.font.size >= Pt(24):
                                run.font.color.rgb = h(pc)
                            else:
                                run.font.color.rgb = h(tc)
                        changes += 1

    output = io.BytesIO(); prs.save(output); output.seek(0)
    return output, changes


# ==================== UI ====================

st.title("🪄 PPT 工具箱")
st.caption("AI 智能生成 · 一键美化 · 通用去水印 · 布局修复")

# ---- 侧边栏：权限管理 ----
with st.sidebar:
    st.subheader("🔐 权限管理")

    if not st.session_state.ai_enabled:
        auth_method = st.radio("选择验证方式", ["🔑 管理密码", "🎫 激活码"], horizontal=True)

        if auth_method == "🔑 管理密码":
            pwd = st.text_input("输入管理密码", type="password")
            if st.button("验证密码", use_container_width=True):
                if pwd == ADMIN_PASSWORD:
                    st.session_state.ai_enabled = True
                    st.session_state.auth_type = "admin"
                    st.success("✅ 管理员验证成功！")
                    st.rerun()
                else:
                    st.error("❌ 密码错误")

        else:
            code = st.text_input("输入激活码", placeholder="PPT-XXXX-XXXX")
            if st.button("激活", use_container_width=True):
                entry = validate_code(code)
                if entry:
                    st.session_state.ai_enabled = True
                    st.session_state.auth_type = "activation"
                    st.session_state.code_used = code
                    mark_code_used(code)
                    st.success("✅ 激活成功！")
                    st.rerun()
                else:
                    st.error("❌ 激活码无效或已被使用")
    else:
        if st.session_state.auth_type == "admin":
            st.success("🔓 管理员模式")
        else:
            st.success(f"🔓 已激活: {st.session_state.code_used}")
        st.caption(f"剩余激活码: {sum(1 for c in VALID_CODES if not c['used'])} 个")

    st.divider()
    st.caption(f"API 余额由管理员提供 · DeepSeek")

# ---- 主界面 Tabs ----
tab1, tab2, tab3, tab4 = st.tabs(["🤖 AI 生成 PPT", "✨ 美化 PPT", "🔓 通用去水印", "🔧 布局修复"])

# ======== Tab 1: AI 生成 PPT 内容（主输出 Kimi 提示词 + 草稿 PPTX） ========
with tab1:
    st.subheader("🤖 AI 策划 PPT 内容 → 复制到 Kimi 生成精美 PPT")

    col1, col2, col3 = st.columns(3)
    with col1:
        theme = st.selectbox("🎨 风格", list(THEMES.keys()))
    with col2:
        slide_count = st.slider("📊 页数", 4, 15, 8)
    with col3:
        mode = "AI 智能策划" if st.session_state.ai_enabled else "模板生成（免费）"
        st.info(f"📋 {mode}")

    prompt = st.text_area(
        "📝 输入主题或详细描述",
        placeholder="短主题：新能源汽车2025市场分析\n\n也可以粘贴长描述：\n介绍当前新能源汽车市场的整体情况，分析特斯拉、比亚迪、蔚来等主要玩家的竞争格局，重点解读2025年电池技术突破对行业的影响...",
        height=130,
    )

    if st.button("🪄 生成 PPT 内容", type="primary", use_container_width=True):
        if not prompt.strip():
            st.warning("请输入内容")
        else:
            if st.session_state.ai_enabled:
                with st.spinner("🔍 第1步：AI 研究市场最佳实践..."):
                    try:
                        data, research, kimi_md = generate_ppt_content(prompt.strip(), theme, slide_count)

                        # === 主输出：Kimi 提示词 ===
                        st.success(f"✅ 内容生成成功！")
                        st.markdown("---")
                        st.subheader("📋 Kimi 提示词（复制下面内容，粘贴到 Kimi 即可生成精美 PPT）")

                        # 用 code block 展示方便复制
                        st.code(kimi_md, language="markdown", line_numbers=False)

                        # 一键复制按钮
                        b64 = base64.b64encode(kimi_md.encode()).decode()
                        st.markdown(
                            f'<a href="data:text/plain;charset=utf-8;base64,{b64}" '
                            f'download="kimi_ppt_prompt.md" '
                            f'style="display:inline-block;padding:8px 20px;background:#6C63FF;color:white;'
                            f'text-decoration:none;border-radius:6px;font-weight:bold;">'
                            f'📥 下载 Markdown 文件</a>'
                            f'&nbsp;&nbsp;<span style="color:#888;font-size:13px;">↑ 下载后用 Kimi 打开或复制粘贴</span>',
                            unsafe_allow_html=True,
                        )

                        # 研究摘要
                        with st.expander("📊 AI 研究摘要"):
                            st.caption(f"**市场背景**: {research.get('market_context', '—')}")
                            st.caption(f"**目标听众**: {research.get('target_audience', '—')}")
                            st.caption(f"**关键角度**: {', '.join(research.get('key_angles', []))}")
                            if research.get("outline"):
                                for o in research["outline"]:
                                    st.caption(f"- **{o.get('section', '')}**: {o.get('focus', '')}")

                        # === 次要输出：PPTX 草稿下载 ===
                        st.markdown("---")
                        st.caption("⬇️ 也可以下载 PPTX 草稿（本地生成，布局较简单；推荐用上面的 Kimi 提示词获得更好效果）")
                        pptx_data = build_pptx(data, theme, slide_count)
                        safe_name = re.sub(r'[^\w一-鿿]', '_', prompt.strip()[:30])
                        st.download_button("📥 下载 PPTX 草稿", pptx_data,
                                           file_name=f"{safe_name}_{theme}.pptx",
                                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                           use_container_width=True)

                    except Exception as e:
                        st.error(f"AI 生成失败: {e}")
                        st.info("自动切换到模板模式...")
                        outline = OUTLINE_LIBRARY.get(theme, OUTLINE_LIBRARY["专业商务"])
                        data = {
                            "title": prompt.strip()[:60],
                            "subtitle": f"{theme} · 智能生成",
                            "slides": [{"title": item, "layout": "bullets",
                                        "bullets": [f"{item}相关内容", "核心要点分析", "实践案例解读"], "notes": ""}
                                       for item in outline[:slide_count - 1]],
                        }
                        pptx_data = build_pptx(data, theme, slide_count)
                        safe_name = re.sub(r'[^\w一-鿿]', '_', prompt.strip()[:30])
                        st.download_button("📥 下载 PPTX 草稿", pptx_data,
                                           file_name=f"{safe_name}_{theme}.pptx",
                                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                           use_container_width=True)
            else:
                st.warning("⚠️ AI 策划需要激活码或管理密码，请在左侧边栏解锁。当前使用模板生成。")
                with st.spinner("📋 正在构建 PPT..."):
                    outline = OUTLINE_LIBRARY.get(theme, OUTLINE_LIBRARY["专业商务"])
                    data = {
                        "title": prompt.strip()[:60],
                        "subtitle": f"{theme} · 智能生成",
                        "slides": [{"title": item, "layout": "bullets",
                                    "bullets": [f"{item} - 核心内容", f"{item} - 关键分析", f"{item} - 实例说明"],
                                    "notes": ""} for item in outline[:slide_count - 1]],
                    }
                    st.success(f"✅ 模板生成成功，共 {len(data['slides']) + 1} 页")
                pptx_data = build_pptx(data, theme, slide_count)
                safe_name = re.sub(r'[^\w一-鿿]', '_', prompt.strip()[:30])
                st.download_button("📥 下载 PPTX", pptx_data,
                                   file_name=f"{safe_name}_{theme}.pptx",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True)

# ======== Tab 2: 美化 ========
with tab2:
    st.subheader("✨ PPT 一键美化")
    c1, c2 = st.columns(2)
    with c1:
        b_theme = st.selectbox("🎨 目标风格", list(THEMES.keys()), key="b_theme")
    with c2:
        intensity = st.selectbox("🖌️ 美化程度", ["轻度（字体）", "中度（字体+配色）", "深度（全部重排）"])

    uploaded_b = st.file_uploader("上传 PPTX", type=["pptx"], key="b_upload")
    if uploaded_b:
        if st.button("✨ 开始美化", type="primary", use_container_width=True):
            with st.spinner("美化中..."):
                result, changes = beautify_pptx(uploaded_b.getvalue(), uploaded_b.name, b_theme, intensity)
                st.success(f"✅ 处理了 {changes} 处文本")
                st.download_button("📥 下载美化版", result,
                                   file_name=f"美化_{b_theme}_{uploaded_b.name}",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True)

# ======== Tab 3: 去水印 ========
with tab3:
    st.subheader("🔓 通用去水印")
    wm = st.text_input("🔍 指定水印文字（可选）", placeholder="留空自动检测 80+ 常见水印")
    uploaded_w = st.file_uploader("上传 PPTX", type=["pptx"], key="w_upload")
    if uploaded_w:
        if st.button("🔓 去水印", use_container_width=True):
            with st.spinner("搜索中..."):
                result, count, found = remove_watermark(uploaded_w.getvalue(), uploaded_w.name,
                                                        wm.strip() if wm.strip() else None)
                if count:
                    st.success(f"✅ 移除 {count} 个标记 ({', '.join(found)})")
                    st.download_button("📥 下载", result, file_name=f"去水印_{uploaded_w.name}",
                                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                       use_container_width=True)
                else:
                    st.info("未检测到水印")

# ======== Tab 4: 修复 ========
with tab4:
    st.subheader("🔧 布局修复")
    uploaded_f = st.file_uploader("上传 PPTX", type=["pptx"], key="f_upload")
    if uploaded_f:
        if st.button("🔧 修复", use_container_width=True):
            with st.spinner("修复中..."):
                result, fixed = fix_layout(uploaded_f.getvalue(), uploaded_f.name)
                st.success(f"✅ 修复 {fixed} 处")
                st.download_button("📥 下载修复版", result, file_name=f"修复_{uploaded_f.name}",
                                   mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                   use_container_width=True)
